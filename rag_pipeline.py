import os
import requests
from dotenv import load_dotenv
from pathlib import Path
import PyPDF2
import docx
import pptx
import re
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter

# Load environment variable
load_dotenv()

# Hugging Face API Setup for LLM
HF_API_URL = "https://api-inference.huggingface.co/models/meta-llama/Meta-Llama-3-8B-Instruct"
HF_API_TOKEN = os.getenv("HF_API_TOKEN")
headers = {"Authorization": f"Bearer {HF_API_TOKEN}"}

# Initialize embeddings for vector database
embedding_model = HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

# File Processing Functions
def extract_text_from_pdf(file_path):
    """Extract text from PDF files."""
    text = ""
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file_path):
    """Extract text from DOCX files."""
    doc = docx.Document(file_path)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint files."""
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_txt(file_path):
    """Extract text from TXT files."""
    with open(file_path, "r", encoding="utf-8") as file:
        return file.read()

def process_course_materials(directory_path):
    """Process all supported files in the directory and extract text."""
    all_text = ""
    file_types = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".pptx": extract_text_from_pptx,
        ".txt": extract_text_from_txt
    }
    
    for file_path in Path(directory_path).glob("**/*"):
        if file_path.suffix.lower() in file_types:
            print(f"Processing {file_path}")
            try:
                text = file_types[file_path.suffix.lower()](file_path)
                all_text += f"\n\n--- Document: {file_path.name} ---\n\n{text}"
            except Exception as e:
                print(f"Error processing {file_path}: {e}")
    
    return all_text

def build_vector_database(text):
    """Build a vector database from the text for semantic search."""
    # Create a text splitter to divide content into smaller chunks
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200
    )
    chunks = text_splitter.split_text(text)
    
    # Create vector store
    vector_store = FAISS.from_texts(chunks, embedding_model)
    
    # Save the vector store to disk (optional)
    vector_store.save_local("faiss_index")
    
    return vector_store

def load_vector_database():
    """Load the vector database from disk."""
    vector_store = FAISS.load_local("faiss_index", embedding_model)
    return vector_store

def analyze_question(question):
    """
    Analyze the question to determine its intent and extract key entities.
    This is a basic implementation - could be enhanced with more sophisticated NLP.
    """
    # Simple keyword extraction
    keywords = re.findall(r'\b[A-Za-z][A-Za-z-]+\b', question.lower())
    # Remove common stop words
    stop_words = {'what', 'when', 'where', 'who', 'why', 'how', 'is', 'are', 'the', 'a', 'an', 'in', 'on', 'at', 'to', 'for', 'with', 'by', 'about', 'like', 'through', 'over', 'before', 'after', 'between', 'under', 'during', 'since', 'without'}
    keywords = [word for word in keywords if word not in stop_words and len(word) > 3]
    
    # Question classification (very simple approach)
    question_type = None
    if question.lower().startswith('what'):
        question_type = 'definition'
    elif question.lower().startswith('how'):
        question_type = 'procedure'
    elif question.lower().startswith('why'):
        question_type = 'explanation'
    elif question.lower().startswith('when'):
        question_type = 'time'
    elif question.lower().startswith('where'):
        question_type = 'location'
    elif question.lower().startswith('who'):
        question_type = 'person'
    else:
        question_type = 'general'
    
    return {
        'keywords': keywords,
        'question_type': question_type,
        'original_question': question
    }

def retrieve_relevant_contexts(vector_store, question, k=3):
    """
    Retrieve the most relevant text passages for the question.
    
    Args:
        vector_store: The FAISS vector store containing the document embeddings
        question: The user's question
        k: Number of documents to retrieve
        
    Returns:
        List of retrieved documents
    """
    # Analyze the question
    analysis = analyze_question(question)
    
    # Get the k most similar contexts from the vector store
    documents = vector_store.similarity_search(
        question, 
        k=k
    )
    
    # Extract the text from the documents
    contexts = [doc.page_content for doc in documents]
    
    return contexts, analysis

def query_llm(prompt):
    """
    Query the language model (Llama 3) to generate a response.
    
    Args:
        prompt: The formatted prompt for the model
        
    Returns:
        The model's response
    """
    try:
        response = requests.post(
            HF_API_URL, 
            headers=headers, 
            json={"inputs": prompt}
        )
        
        if response.status_code == 200:
            api_response = response.json()
            if api_response and isinstance(api_response, list) and "generated_text" in api_response[0]:
                generated_text = api_response[0]["generated_text"].strip()
                
                # Remove the prompt from the beginning if present
                if generated_text.startswith(prompt):
                    generated_text = generated_text[len(prompt):].strip()
                
                return generated_text
            else:
                print("Error: Invalid API response format.")
                return "Error: Could not generate a valid response."
        else:
            print(f"Error {response.status_code}: {response.text}")
            return f"Error {response.status_code}: Could not get a response from the model."
    except Exception as e:
        print(f"Exception during API call: {e}")
        return f"Error: {str(e)}"

def format_rag_prompt(question, contexts):
    """
    Format the RAG prompt by combining the question and retrieved contexts.
    
    Args:
        question: The user's question
        contexts: List of retrieved document contexts
        
    Returns:
        Formatted prompt for the language model
    """
    context_text = "\n\n".join([f"Context {i+1}:\n{context}" for i, context in enumerate(contexts)])
    
    prompt = f"""You are an educational assistant that helps answer questions based on course materials.
Use ONLY the following contexts from course materials to answer the question. If the answer is not contained in the contexts, say "I don't have enough information to answer this question based on the course materials."

{context_text}

Question: {question}

Answer: """
    
    return prompt

def generate_answer(question, contexts, analysis):
    """
    Generate an answer based on the retrieved contexts.
    
    Args:
        question: The user's question
        contexts: List of retrieved document contexts
        analysis: Analysis of the question
        
    Returns:
        Generated answer
    """
    prompt = format_rag_prompt(question, contexts)
    answer = query_llm(prompt)
    
    # Post-processing of the answer (optional)
    # You could implement additional filtering, fact-checking, etc. here
    
    return answer

def rag_pipeline(question, vector_store=None):
    """
    The complete RAG pipeline:
    1. Analyze the question
    2. Retrieve relevant contexts
    3. Generate answer using LLM with retrieved contexts
    
    Args:
        question: The user's question
        vector_store: Optional pre-loaded vector store
        
    Returns:
        Generated answer
    """
    # Load vector store if not provided
    if vector_store is None:
        try:
            vector_store = load_vector_database()
        except:
            print("Error: Vector database not found. Please build it first.")
            return "Error: Knowledge base not available."
    
    # Retrieve relevant contexts
    contexts, analysis = retrieve_relevant_contexts(vector_store, question)
    
    # Generate answer
    answer = generate_answer(question, contexts, analysis)
    
    return {
        'question': question,
        'answer': answer,
        'contexts': contexts,
        'analysis': analysis
    }

def main():
    """Main function to demonstrate RAG pipeline."""
    # Check if vector database exists
    if not Path("faiss_index").exists():
        print("Building vector database from course materials...")
        course_directory = "course_materials"
        course_text = process_course_materials(course_directory)
        
        # Save extracted text to a file (optional)
        with open("extracted_course_content.txt", "w", encoding="utf-8") as f:
            f.write(course_text)
        print("Extracted text saved to extracted_course_content.txt")
        
        vector_store = build_vector_database(course_text)
        print("Vector database built and saved to 'faiss_index'")
    else:
        print("Loading existing vector database...")
        vector_store = load_vector_database()
    
    # Interactive question answering
    print("\n=== Course Content Question Answering System ===")
    print("Type 'exit' to quit\n")
    
    while True:
        question = input("\nEnter your question: ")
        if question.lower() == 'exit':
            break
        
        print("\nProcessing your question...\n")
        result = rag_pipeline(question, vector_store)
        
        print("Answer:", result['answer'])
        print("\n---")
        print("Question analysis:", result['analysis']['question_type'])
        print("Keywords identified:", ", ".join(result['analysis']['keywords']))
        print("---\n")
        # Uncomment to show the retrieved contexts
        # print("Contexts used:")
        # for i, context in enumerate(result['contexts']):
        #     print(f"Context {i+1}:\n{context[:150]}...\n")

if __name__ == "__main__":
    main()